# -*- coding: utf-8 -*-
"""
Generates Fitness_Synergy_Presentation.pptx
Run:  python build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Theme ----------
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
LBLUE  = RGBColor(0xE8, 0xF0, 0xFE)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1F, 0x2A, 0x37)
GRAY   = RGBColor(0x5B, 0x66, 0x72)
LGRAY  = RGBColor(0xF3, 0xF5, 0xF8)

FONT = "Calibri"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def box(slide, x, y, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    return sp


def rounded(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for seg in para:
            txt, size, color, bold, italic = (list(seg) + [False, False])[:5]
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, italic)
    return tb


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


def bullets(slide, x, y, w, h, items, size=16, gap=10, color=DARK, bullet_color=BLUE):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        # it: (text, level) or (text, level, bold)
        txt = it[0]; level = it[1] if len(it) > 1 else 0
        bold = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        mark = "▸ " if level == 0 else "•  "
        r = p.add_run(); r.text = mark; _set(r, size, bullet_color if level == 0 else ORANGE, True)
        r2 = p.add_run(); r2.text = txt
        _set(r2, size if level == 0 else size-1, color, bold)
        if level == 1:
            p.level = 1
    return tb


# ---------- Slide builders ----------
def header(slide, kicker, title, tag=None):
    box(slide, 0, 0, EMU_W, Inches(1.25), fill=NAVY)
    box(slide, 0, Inches(1.25), EMU_W, Inches(0.06), fill=ORANGE)
    text(slide, Inches(0.6), Inches(0.16), Inches(11), Inches(0.4),
         [[(kicker, 13, ORANGE, True)]])
    text(slide, Inches(0.6), Inches(0.5), Inches(11.5), Inches(0.7),
         [[(title, 26, WHITE, True)]])
    if tag:
        text(slide, Inches(10.4), Inches(0.4), Inches(2.6), Inches(0.6),
             [[(tag, 12, LBLUE, False)]], align=PP_ALIGN.RIGHT)


CODE_BG   = RGBColor(0x0E, 0x14, 0x22)
CODE_TXT  = RGBColor(0xE6, 0xED, 0xF3)
CODE_CMT  = RGBColor(0x7E, 0xE7, 0x87)
CODE_KW   = RGBColor(0xF5, 0x9E, 0x0B)


def code_slide(kicker, title, code_lines, callouts, accent=BLUE, tag=None):
    s = prs.slides.add_slide(BLANK)
    header(s, kicker, title, tag=tag)
    # left: code panel
    panel = rounded(s, Inches(0.5), Inches(1.55), Inches(7.0), Inches(5.4), fill=CODE_BG)
    box(s, Inches(0.5), Inches(1.55), Inches(7.0), Inches(0.42), fill=RGBColor(0x1B,0x26,0x3B))
    text(s, Inches(0.75), Inches(1.6), Inches(6.6), Inches(0.35),
         [[("● ● ●   backend code", 11, RGBColor(0x8B,0x98,0xA8), False, False, "Consolas")]])
    tb = s.shapes.add_textbox(Inches(0.72), Inches(2.12), Inches(6.65), Inches(4.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2); p.line_spacing = 1.0
        st = ln.strip()
        color = CODE_CMT if (st.startswith("//") or st.startswith("--")) else CODE_TXT
        r = p.add_run(); r.text = ln if ln else " "
        _set(r, 11.5, color, False, st.startswith("//") or st.startswith("--"), "Consolas")
    # right: callouts
    text(s, Inches(7.8), Inches(1.65), Inches(5.0), Inches(0.5),
         [[("Ano'ng nangyayari", 17, NAVY, True)]])
    yy = Inches(2.3)
    for i, (c) in enumerate(callouts):
        h = Inches(1.02)
        rounded(s, Inches(7.8), yy, Inches(4.95), h, fill=LGRAY)
        box(s, Inches(7.8), yy, Inches(0.55), h, fill=accent)
        text(s, Inches(7.8), yy+Inches(0.28), Inches(0.55), Inches(0.5),
             [[(str(i+1), 20, WHITE, True)]], align=PP_ALIGN.CENTER)
        text(s, Inches(8.5), yy+Inches(0.12), Inches(4.15), h-Inches(0.2),
             [[(c, 12.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + h + Inches(0.12)
    return s


def divider(num, who, title, subtitle, accent=BLUE):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
    box(s, 0, Inches(3.05), EMU_W, Inches(0.06), fill=ORANGE)
    text(s, Inches(0.9), Inches(1.5), Inches(3), Inches(1.6),
         [[(num, 96, accent, True)]])
    text(s, Inches(3.2), Inches(1.7), Inches(9), Inches(0.6),
         [[("PART " + num + "  ·  " + who, 16, ORANGE, True)]])
    text(s, Inches(3.2), Inches(2.15), Inches(9.2), Inches(1.1),
         [[(title, 34, WHITE, True)]])
    text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2),
         [[(subtitle, 17, LBLUE, False)]])
    return s


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
box(s, 0, Inches(4.7), EMU_W, Inches(0.08), fill=ORANGE)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
     [[("DATABASE MANAGEMENT  ·  SYSTEM DEFENSE", 15, ORANGE, True)]])
text(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(1.6),
     [[("Fitness Synergy", 60, WHITE, True)],
      [("Gym Management System", 30, LBLUE, False)]])
text(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.8),
     [[("A web application for Fitness Synergy Lipa — built on a relational MySQL database.", 16, WHITE, False)],
      [("Team M.G.A.B. Co.   |   Lead: Marco Andrei R. Belen", 15, ORANGE, True)],
      [("React + Vite  →  PHP REST API  →  MySQL (MariaDB)", 13, GRAY, False)]])
notes(s, "Marco opens. Greet, introduce the group and the project, say each member will present a slice of the database. Keep calm, this is the easiest slide.")

# ============================================================
# SLIDE 2 — MEET THE TEAM / AGENDA
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "AGENDA", "Who presents what")
cards = [
    ("1", "Foundation", "Architecture & Database Design", "Marco Andrei Belen (Lead)", BLUE),
    ("2", "Getting Data In", "Membership Lifecycle — writes & transactions", "Ariz Luiz Guillarte", ORANGE),
    ("3", "Daily Operations", "Attendance & Reference Data", "Brix Martin Paña", BLUE),
    ("4", "Insight & Trust", "Reporting + Security & Integrity", "Gabriel Caporado", ORANGE),
]
cx = Inches(0.6); cw = Inches(5.95); ch = Inches(2.35); gap = Inches(0.25)
positions = [(Inches(0.6), Inches(1.65)), (Inches(6.78), Inches(1.65)),
             (Inches(0.6), Inches(4.3)), (Inches(6.78), Inches(4.3))]
for (num, t, sub, who, ac), (x, y) in zip(cards, positions):
    rounded(s, x, y, cw, ch, fill=LGRAY)
    box(s, x, y, Inches(0.12), ch, fill=ac)
    text(s, x+Inches(0.4), y+Inches(0.2), Inches(1.2), Inches(1),
         [[(num, 44, ac, True)]])
    text(s, x+Inches(1.55), y+Inches(0.28), Inches(4.2), Inches(1.8),
         [[(t, 19, NAVY, True)],
          [(sub, 13.5, DARK, False)],
          [("▶  " + who, 13, ac, True)]], space_after=7)
notes(s, "The presentation follows the life of the data: design it, get it in, operate on it, turn it into insight. Each member owns one stage. The demo is one continuous story.")

# ============================================================
# SLIDE 3 — WHAT IS FITNESS SYNERGY
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "OVERVIEW", "What is Fitness Synergy?")
text(s, Inches(0.6), Inches(1.6), Inches(7.2), Inches(2),
     [[("A complete gym management system that replaces pen-and-paper logbooks for ",16,DARK),
       ("Fitness Synergy Lipa",16,NAVY,True),(".",16,DARK)],
      [("Front-desk staff use it every day to register members, accept payments, "
        "track attendance, and see how the business is doing.",16,DARK)]], space_after=10)
feat = [
    "Member & walk-in registration", "Membership renewals & installments",
    "Daily attendance / time-in", "Plans, promos & expense tracking",
    "Revenue & sales reporting", "Secure login + full audit trail",
]
rounded(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.9), fill=LBLUE)
text(s, Inches(0.9), Inches(4.2), Inches(11), Inches(0.5),
     [[("What the system does", 17, NAVY, True)]])
col1 = feat[:3]; col2 = feat[3:]
bullets(s, Inches(0.95), Inches(4.85), Inches(6), Inches(2),
        [(f, 0) for f in col1], size=15, gap=12)
bullets(s, Inches(6.9), Inches(4.85), Inches(5.7), Inches(2),
        [(f, 0) for f in col2], size=15, gap=12)
notes(s, "Real client in Lipa City. Emphasize: this is a real, working system on a live database, not a mockup.")

# ============================================================
# SLIDE 4 — 3-TIER ARCHITECTURE
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "ARCHITECTURE", "Three-tier design")
tiers = [
    ("FRONTEND", "React + Vite", "Single-page app in the browser.\nButtons, forms, tables.\nNever touches the database.", BLUE),
    ("BACKEND", "PHP REST API", "~55 endpoint files.\nChecks login & rules,\nruns the SQL.", ORANGE),
    ("DATABASE", "MySQL / MariaDB", "13 tables.\nThe only thing the\nbackend talks to.", NAVY),
]
x = Inches(0.6); w = Inches(3.7); y = Inches(1.85); h = Inches(2.9); gapx = Inches(0.62)
for i, (cap, tech, desc, ac) in enumerate(tiers):
    cx = x + i*(w+gapx)
    rounded(s, cx, y, w, h, fill=LGRAY, line=ac)
    box(s, cx, y, w, Inches(0.7), fill=ac)
    text(s, cx, y+Inches(0.08), w, Inches(0.55),
         [[(cap, 16, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, cx, y+Inches(0.85), w, Inches(0.5),
         [[(tech, 18, NAVY, True)]], align=PP_ALIGN.CENTER)
    text(s, cx+Inches(0.25), y+Inches(1.5), w-Inches(0.5), Inches(1.3),
         [[(line, 13, DARK, False)] for line in desc.split("\n")],
         align=PP_ALIGN.CENTER, space_after=2)
    if i < 2:
        text(s, cx+w-Inches(0.05), y+Inches(1.05), Inches(0.7), Inches(0.6),
             [[("→", 30, ORANGE, True)]], align=PP_ALIGN.CENTER)
rounded(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(1.7), fill=NAVY)
text(s, Inches(0.95), Inches(5.35), Inches(11.4), Inches(1.4),
     [[("Why this matters:  ", 16, ORANGE, True),
       ("the browser has no direct line to the database. Every request is forced through "
        "the backend, which checks your login and validates input first. ", 16, WHITE, False)],
      [("If the frontend could run SQL, anyone could open the dev console and wipe the data.",
        15, LBLUE, True)]], space_after=8)
notes(s, "Restaurant analogy: frontend = dining room, backend = kitchen, database = pantry. The browser never touches the DB directly.")

# ============================================================
# SLIDE 5 — REQUEST LIFECYCLE
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "ARCHITECTURE", "One click, end to end")
steps = [
    ("1", "Click 'Time In'", "Browser (React)"),
    ("2", "Send request + token", "api.js  →  HTTP/JSON"),
    ("3", "Check rules", "time_in.php"),
    ("4", "Run SQL", "INSERT INTO attendance"),
    ("5", "Save row", "MySQL"),
    ("6", "Return JSON", "{ success: true }"),
]
y = Inches(2.1); w = Inches(1.92); h = Inches(2.2); x0 = Inches(0.55); gapx = Inches(0.12)
for i, (n, t, sub) in enumerate(steps):
    cx = x0 + i*(w+gapx)
    ac = BLUE if i % 2 == 0 else ORANGE
    rounded(s, cx, y, w, h, fill=LGRAY)
    box(s, cx, y, w, Inches(0.55), fill=ac)
    text(s, cx, y+Inches(0.04), w, Inches(0.5), [[("STEP "+n, 11, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, cx+Inches(0.1), y+Inches(0.75), w-Inches(0.2), Inches(0.9),
         [[(t, 14, NAVY, True)]], align=PP_ALIGN.CENTER)
    text(s, cx+Inches(0.1), y+Inches(1.45), w-Inches(0.2), Inches(0.7),
         [[(sub, 11, GRAY, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(1.5),
     [[("A button click becomes one JavaScript call → one PHP file → one SQL statement → one row.",
        17, DARK, True)],
      [("That is the entire system in miniature — and the pattern every feature follows (a REST API).",
        15, GRAY, False)]], space_after=8)
notes(s, "Trace one click live: api.js -> time_in.php -> SQL -> JSON. One-liner: one click -> one PHP file -> one SQL statement -> one row.")

# ============================================================
# DIVIDER PART 1
# ============================================================
divider("1", "Marco (Leader)", "Foundation: Architecture & Database Design",
        "The blueprint everyone builds on — the 13 tables, how they relate, and the engineering choices behind them.")

# ============================================================
# SLIDE 6 — THE DATABASE (13 TABLES)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DATABASE DESIGN", "13 tables, grouped by purpose", tag="MariaDB 11.4 · InnoDB")
groups = [
    ("CORE DATA", ["members", "payments", "attendance"], BLUE),
    ("REFERENCE", ["plans", "promos", "expenses"], ORANGE),
    ("REPORTING", ["monthly_targets", "bank_deposits"], BLUE),
    ("SECURITY & SYSTEM", ["admins", "sessions", "login_attempts", "activity_log", "contract_counter"], NAVY),
]
positions = [(Inches(0.6), Inches(1.7)), (Inches(6.78), Inches(1.7)),
             (Inches(0.6), Inches(4.35)), (Inches(6.78), Inches(4.35))]
cw = Inches(5.95); ch = Inches(2.3)
for (cap, tbls, ac), (x, y) in zip(groups, positions):
    rounded(s, x, y, cw, ch, fill=LGRAY)
    box(s, x, y, cw, Inches(0.55), fill=ac)
    text(s, x+Inches(0.3), y+Inches(0.04), cw, Inches(0.5), [[(cap, 14, WHITE, True)]])
    chips = "   ".join(tbls)
    text(s, x+Inches(0.3), y+Inches(0.75), cw-Inches(0.6), Inches(1.5),
         [[(t, 15, NAVY, True)] for t in tbls] if len(tbls) <= 3 else
         [[(tbls[0], 15, NAVY, True), ("   "+tbls[1], 15, NAVY, True)],
          [(tbls[2], 15, NAVY, True), ("   "+tbls[3], 15, NAVY, True)],
          [(tbls[4], 15, NAVY, True)]],
         space_after=6)
notes(s, "Point to the four groups, don't read every table.")

# ============================================================
# SLIDE 7 — ERD / RELATIONSHIPS
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DATABASE DESIGN", "How the core tables relate (ERD)")
# plans
rounded(s, Inches(5.2), Inches(1.7), Inches(2.9), Inches(1.0), fill=ORANGE)
text(s, Inches(5.2), Inches(1.85), Inches(2.9), Inches(0.8),
     [[("plans", 18, WHITE, True)],[("plan_id (PK)", 11, LBLUE, False)]], align=PP_ALIGN.CENTER, space_after=0)
# members
rounded(s, Inches(5.2), Inches(3.35), Inches(2.9), Inches(1.0), fill=BLUE)
text(s, Inches(5.2), Inches(3.5), Inches(2.9), Inches(0.8),
     [[("members", 18, WHITE, True)],[("member_id (PK)", 11, LBLUE, False)]], align=PP_ALIGN.CENTER, space_after=0)
# payments
rounded(s, Inches(1.2), Inches(5.0), Inches(2.9), Inches(1.0), fill=NAVY)
text(s, Inches(1.2), Inches(5.15), Inches(2.9), Inches(0.8),
     [[("payments", 18, WHITE, True)],[("member_id FK · plan_id FK", 10, LBLUE, False)]], align=PP_ALIGN.CENTER, space_after=0)
# attendance
rounded(s, Inches(9.2), Inches(5.0), Inches(2.9), Inches(1.0), fill=NAVY)
text(s, Inches(9.2), Inches(5.15), Inches(2.9), Inches(0.8),
     [[("attendance", 18, WHITE, True)],[("member_id FK", 10, LBLUE, False)]], align=PP_ALIGN.CENTER, space_after=0)
# connectors (simple lines via thin boxes)
def connector(x1, y1, x2, y2, label, lx, ly):
    ln = s.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = GRAY; ln.line.width = Pt(2)
    if label:
        text(s, lx, ly, Inches(2.2), Inches(0.4), [[(label, 11, ORANGE, True)]], align=PP_ALIGN.CENTER)
connector(Inches(6.65), Inches(2.7), Inches(6.65), Inches(3.35), "1 → many", Inches(6.8), Inches(2.85))
connector(Inches(5.6), Inches(4.35), Inches(2.9), Inches(5.0), "makes", Inches(3.6), Inches(4.5))
connector(Inches(7.7), Inches(4.35), Inches(10.4), Inches(5.0), "checks in", Inches(8.4), Inches(4.5))
text(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(1),
     [[("One plan → many members  ·  one member → many payments & attendance rows.  ", 14, DARK, True),
       ("A walk-in is a payment with member_id = NULL.", 14, ORANGE, True)]])
notes(s, "Say each relationship out loud; the walk-in (member_id NULL) is the highlight.")

# ============================================================
# SLIDE 8 — WHY IT'S WELL DESIGNED
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DATABASE DESIGN", "Engineering decisions we can defend")
items = [
    ("Third Normal Form (3NF)", "Plan price & duration live only in plans; members store a plan_id foreign key. No duplicated data, no update anomalies."),
    ("InnoDB, not MyISAM", "Gives us transactions (all-or-nothing writes) and foreign keys. We deliberately converted promos & expenses."),
    ("Foreign keys with delete rules", "members.plan_id is RESTRICT (can't delete a plan in use); member deletes CASCADE to their payments & attendance."),
    ("Indexes on hot columns", "payment_date, expiration_date, member_id, plan_id — so reports don't full-scan the tables."),
    ("Deliberate snapshots", "payments.amount is copied at sale time, so changing a plan's price later never rewrites old receipts."),
    ("utf8mb4 charset", "Stores ₱ and — correctly. A real bug we found and fixed."),
]
y = Inches(1.65)
for i, (t, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = Inches(0.6) + col*Inches(6.2)
    yy = y + row*Inches(1.7)
    rounded(s, x, yy, Inches(5.95), Inches(1.55), fill=LGRAY)
    box(s, x, yy, Inches(0.1), Inches(1.55), fill=BLUE if col == 0 else ORANGE)
    text(s, x+Inches(0.32), yy+Inches(0.13), Inches(5.5), Inches(1.4),
         [[(t, 14.5, NAVY, True)],[(d, 11.5, DARK, False)]], space_after=4)
notes(s, "If rushed, emphasize 3NF and InnoDB/transactions. The next slides show the actual code behind these.")

# ============================================================
# SLIDE 8c — DB CONNECTION (MARCO)
# ============================================================
_m1 = code_slide(
    "CODE · PART 1", "How we connect: db.php",
    [
        "// Load secret credentials from .env (not in code)",
        "$host = $env['DB_HOST'];  $dbname = $env['DB_NAME'];",
        "$user = $env['DB_USER'];  $pass   = $env['DB_PASS'];",
        "",
        "// Open ONE connection through PDO",
        "$conn = new PDO(",
        '  "mysql:host=$host;dbname=$dbname;charset=utf8mb4",',
        "  $user, $pass,",
        "  [",
        "    PDO::ATTR_ERRMODE => PDO::ERRMODE_EXCEPTION,",
        "    PDO::ATTR_TIMEOUT => 5,",
        "  ]",
        ");",
        "$conn->exec(\"SET time_zone = '+08:00'\"); // PH time",
    ],
    [
        "WHERE: credentials load from a separate .env file — never hard-coded, so passwords stay out of the source code.",
        "HOW: PDO opens one connection; charset=utf8mb4 lets it store the peso sign and dash correctly.",
        "ERRMODE_EXCEPTION: database errors are thrown and caught, never failing silently.",
        "Time zone is set to +08:00 so every saved date matches Philippine time.",
    ],
    accent=BLUE, tag="Backend/db.php",
)
notes(_m1, "This is the foundation: every other PHP file includes db.php to get $conn. Mention .env keeps the password secret.")

# ============================================================
# SLIDE 8d — SCHEMA: members TABLE (MARCO)
# ============================================================
_m2 = code_slide(
    "CODE · PART 1", "Schema design: the members table",
    [
        "CREATE TABLE members (",
        "  member_id   INT PRIMARY KEY AUTO_INCREMENT,",
        "  full_name   VARCHAR(50) NOT NULL,",
        "  contact_number VARCHAR(20),",
        "  age         TINYINT UNSIGNED,",
        "  gender      ENUM('Male','Female','Other'),",
        "  contract_id VARCHAR(50) UNIQUE,    -- FS-000001",
        "  plan_id     INT,                   -- FK -> plans",
        "  start_date      DATE,",
        "  expiration_date DATE,",
        "  is_installment    TINYINT(1) DEFAULT 0,",
        "  installment_total DECIMAL(10,2) DEFAULT 0",
        ") ENGINE=InnoDB;",
    ],
    [
        "PRIMARY KEY + AUTO_INCREMENT: a unique member_id is generated automatically for every new member.",
        "Each column uses the correct data type — VARCHAR (text), TINYINT/INT (numbers), DATE, DECIMAL (money).",
        "ENUM('Male','Female','Other'): the database only accepts these exact values, so no invalid data gets in.",
        "UNIQUE on contract_id: the DB itself blocks two members from sharing the same FS number.",
    ],
    accent=BLUE, tag="members table (schema)",
)
notes(_m2, "This is real DDL (Data Definition Language). Point out PK, the data types, ENUM, DEFAULT, and UNIQUE.")

# ============================================================
# SLIDE 8e — KEYS & CONSTRAINTS (MARCO)
# ============================================================
_m3 = code_slide(
    "CODE · PART 1", "Data integrity: keys & constraints",
    [
        "-- A member must belong to a real plan",
        "ALTER TABLE members",
        "  ADD FOREIGN KEY (plan_id)",
        "      REFERENCES plans(plan_id);",
        "",
        "-- Delete a member -> remove their child rows too",
        "ALTER TABLE attendance",
        "  ADD FOREIGN KEY (member_id)",
        "      REFERENCES members(member_id)",
        "      ON DELETE CASCADE;",
        "",
        "ALTER TABLE payments",
        "  ADD FOREIGN KEY (member_id)",
        "      REFERENCES members(member_id)",
        "      ON DELETE CASCADE;",
    ],
    [
        "A FOREIGN KEY links two tables and is enforced by the database itself, not by the app code.",
        "members.plan_id -> plans: you cannot assign a member to a plan that does not exist.",
        "ON DELETE CASCADE: deleting a member automatically deletes their attendance and payment rows.",
        "This is integrity at the database level — the user interface cannot bypass it.",
    ],
    accent=BLUE, tag="Foreign keys / constraints",
)
notes(_m3, "Tie back to the ERD. Note RESTRICT vs CASCADE: plans use RESTRICT (can't delete a plan in use).")

# ============================================================
# SLIDE 8f — INDEXES (MARCO)
# ============================================================
_m4 = code_slide(
    "CODE · PART 1", "Performance: indexes",
    [
        "ALTER TABLE payments",
        "  ADD INDEX idx_payment_date    (payment_date),",
        "  ADD INDEX idx_payments_member (member_id);",
        "",
        "ALTER TABLE members",
        "  ADD INDEX idx_expiration_date (expiration_date),",
        "  ADD UNIQUE KEY (contract_id);",
        "",
        "-- Attendance lookups by member and time",
        "ALTER TABLE attendance",
        "  ADD INDEX idx_attendance_member (member_id),",
        "  ADD INDEX idx_time_in           (time_in);",
    ],
    [
        "An index is like a book index: the DB jumps straight to matching rows instead of reading the whole table.",
        "Without an index, every query does a full table scan — which gets slow as data grows.",
        "We indexed the columns the reports filter by: payment_date, member_id, expiration_date.",
        "A UNIQUE index also enforces a rule — no duplicate contract_id is ever allowed.",
    ],
    accent=BLUE, tag="Indexes (performance)",
)
notes(_m4, "Indexes = speed. Connect to the reports in Part 4: those SUM/GROUP BY queries are fast because of these.")

# ============================================================
# SLIDE 8b — CODE LOGIC (MARCO)
# ============================================================
_marco_code = code_slide(
    "CODE LOGIC · PART 1", "Tracing the logic: time_in.php",
    [
        "// 1. Guard: already timed in today?",
        '$check = $conn->prepare(',
        '   "SELECT log_id FROM attendance',
        '    WHERE member_id = :id',
        "      AND DATE(time_in) = CURRENT_DATE()\");",
        "$check->execute([':id' => $member_id]);",
        "",
        "if ($check->rowCount() > 0) {",
        "    // 2. Yes -> reject (no double entry)",
        '    echo \'{"error":"Already timed in"}\';',
        "} else {",
        "    // 3. No -> insert the record",
        '    $conn->prepare(',
        '      "INSERT INTO attendance (member_id, time_in)',
        '       VALUES (:id, NOW())")',
        "     ->execute([':id' => $member_id]);",
        '    echo \'{"success":true}\';',
        "}",
    ],
    [
        "WHAT: records one gym visit for this member, dated today.",
        "HOW (guard): a SELECT first checks the attendance table for a row with this member_id dated today.",
        "DECISION: if a row already exists -> return an error; if none -> run the INSERT.",
        "WHERE: the new row goes into the attendance table; ':id' is a bound parameter, so it is safe from SQL injection.",
    ],
    accent=BLUE, tag="Backend/time_in.php",
)
notes(_marco_code, "Walk the 4 callouts on the slide: guard SELECT, the decision, the INSERT, and the bound parameter (:id).")

# ============================================================
# DIVIDER PART 2
# ============================================================
divider("2", "Ariz Luiz Guillarte", "Getting Data In: Membership Lifecycle",
        "How a real person becomes a row — registration, walk-ins, renewals and installments. The write path and transactions.",
        accent=ORANGE)

# ============================================================
# SLIDE 9 — REGISTRATION TRANSACTION
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "WRITES & TRANSACTIONS", "Registering a member is atomic")
text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.7),
     [[("One click writes to TWO tables — and both must succeed together.", 16, DARK, True)]])
# flow
flow = [("members", "new member row", BLUE), ("payments", "first payment row", NAVY)]
for i, (t, d, ac) in enumerate(flow):
    x = Inches(0.9) + i*Inches(3.4)
    rounded(s, x, Inches(2.35), Inches(2.9), Inches(1.0), fill=ac)
    text(s, x, Inches(2.5), Inches(2.9), Inches(0.8),
         [[("INSERT → "+t, 15, WHITE, True)],[(d, 11, LBLUE, False)]], align=PP_ALIGN.CENTER, space_after=0)
rounded(s, Inches(7.9), Inches(2.35), Inches(4.8), Inches(1.0), fill=ORANGE)
text(s, Inches(7.9), Inches(2.5), Inches(4.8), Inches(0.8),
     [[("Wrapped in ONE transaction", 15, WHITE, True)],
      [("commit() succeeds · rollBack() undoes all", 11, WHITE, False)]], align=PP_ALIGN.CENTER, space_after=0)
bullets(s, Inches(0.6), Inches(3.9), Inches(12), Inches(3),
        [("Atomic: if the payment fails, the member insert is rolled back — never a half-saved record.", 0),
         ("Safe contract IDs: a locked contract_counter row hands out FS-000001, FS-000002… with no race condition (better than MAX(id)+1).", 0),
         ("SQL-injection proof: every value is a bound parameter in a prepared statement, never glued into the SQL.", 0),
         ("Defense-in-depth validation: contact must be 09xxxxxxxxx, age 1–120, DOB not in the future, duplicate-name guard.", 0)],
        size=15, gap=13)
notes(s, "Highlight slide for Member 2. Open add_member.php and show beginTransaction / commit / rollBack and the contract_counter UPDATE.")

# ============================================================
# SLIDE 10 — WALK-INS
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "WRITES & TRANSACTIONS", "Walk-ins: one table, two purposes")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(4.5),
        [("A guest doesn't need a member record.", 0),
         ("Stored as a payments row with member_id = NULL and customer_type = 'Walk-in'.", 0),
         ("The guest's name & contact live on the payment itself.", 0),
         ("Same-day duplicate guard stops double entries.", 0),
         ("After 7 visits, the system recommends converting them to a member.", 0)],
        size=15, gap=14)
rounded(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.6), fill=LBLUE)
text(s, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.5), [[("Smart touch: walkin_key", 16, NAVY, True)]])
text(s, Inches(7.2), Inches(2.55), Inches(5.2), Inches(3.5),
     [[("A generated STORED column in the database automatically builds a key from the guest's "
        "contact number (or their name if no contact).", 13.5, DARK, False)],
      [("It lets us group repeat walk-ins and count their visits — the database computes it for us, "
        "we never set it by hand.", 13.5, DARK, False)],
      [("This is an advanced relational feature, not just a plain column.", 13.5, ORANGE, True)]], space_after=10)
notes(s, "One source of truth for all incoming money. The walkin_key generated column is a strong point to mention if the prof asks about advanced SQL.")

# ============================================================
# SLIDE 11 — RENEWALS & INSTALLMENTS
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "WRITES & TRANSACTIONS", "Renewals & installment payments")
two = [
    ("Renewals", BLUE, [
        "Extends expiration_date by the plan duration (+ bonus days).",
        "Inserts another payments row — full history is kept.",
        "Guard: can't renew if an installment balance is unpaid.",
    ]),
    ("Installments", ORANGE, [
        "Member pays in parts against installment_total.",
        "Each partial payment is recorded and the balance tracked.",
        "Lets students/seniors spread a 12-month plan over time.",
    ]),
]
for i, (t, ac, its) in enumerate(two):
    x = Inches(0.6) + i*Inches(6.2)
    rounded(s, x, Inches(1.7), Inches(5.95), Inches(4.6), fill=LGRAY)
    box(s, x, Inches(1.7), Inches(5.95), Inches(0.65), fill=ac)
    text(s, x+Inches(0.3), Inches(1.78), Inches(5.5), Inches(0.5), [[(t, 17, WHITE, True)]])
    bullets(s, x+Inches(0.3), Inches(2.6), Inches(5.4), Inches(3.5),
            [(z, 0) for z in its], size=14.5, gap=14, bullet_color=ac)
notes(s, "Renewals/installments are UPDATEs + new payment rows. The installment-balance guard shows business rules enforced in the backend.")

# ============================================================
# SLIDE 11b — CODE LOGIC (MEMBER 2)
# ============================================================
code_slide(
    "CODE LOGIC · PART 2", "The transaction: add_member.php",
    [
        "$conn->beginTransaction();          // start",
        "",
        "// Next contract no. (row is locked)",
        '$conn->prepare("UPDATE contract_counter',
        '   SET last_seq = last_seq + 1 WHERE id=1")->execute();',
        '$seq = $conn->query("SELECT last_seq',
        '   FROM contract_counter WHERE id=1")->fetchColumn();',
        "$contract_id = sprintf('FS-%06d', $seq); // FS-000014",
        "",
        "$insertMember->execute([...]);       // 1) members",
        "$new_id = $conn->lastInsertId();",
        "$insertPayment->execute([            // 2) payments",
        "   ':member_id' => $new_id, ...]);",
        "",
        "$conn->commit();                     // both succeeded",
        "// catch: $conn->rollBack();         // undo everything",
    ],
    [
        "beginTransaction() - everything after this succeeds together (all-or-nothing).",
        "A locked contract_counter row means no two staff can grab the same FS number at the same time.",
        "Two INSERTs: members then payments, linked by lastInsertId().",
        "commit() if both succeed; rollBack() on error -> never a half-saved record.",
    ],
    accent=ORANGE, tag="Backend/add_member.php",
)

# ============================================================
# DIVIDER PART 3
# ============================================================
divider("3", "Brix Martin Paña", "Daily Operations: Attendance & Reference Data",
        "The day-to-day: checking members in, spotting expiring memberships, and managing plans, promos & expenses.")

# ============================================================
# SLIDE 12 — ATTENDANCE
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DAILY OPERATIONS", "Attendance / Time-in")
bullets(s, Inches(0.6), Inches(1.7), Inches(7.0), Inches(4.5),
        [("Member arrives → staff click Time-In.", 0),
         ("Guard query first: SELECT ... WHERE member_id = ? AND DATE(time_in) = CURRENT_DATE() — no double check-ins per day.", 0),
         ("If clear: INSERT INTO attendance (member_id, time_in) VALUES (?, NOW()).", 0),
         ("attendance.member_id is a foreign key with ON DELETE CASCADE — remove a member and their visit history goes too.", 0),
         ("Indexed (member_id, time_in) so the daily list is fast.", 0)],
        size=15, gap=13)
rounded(s, Inches(7.9), Inches(1.7), Inches(4.8), Inches(4.5), fill=NAVY)
text(s, Inches(8.2), Inches(2.0), Inches(4.2), Inches(0.5), [[("Live demo", 16, ORANGE, True)]])
text(s, Inches(8.2), Inches(2.6), Inches(4.2), Inches(3.4),
     [[("Time-in the member that Ariz just registered.", 14, WHITE, True)],
      [("Watch the Live Feed update, then show the new row in the attendance table.", 13.5, LBLUE, False)],
      [("Same record flows from one presenter to the next.", 13.5, ORANGE, True)]], space_after=12)
notes(s, "Use the member just registered. This is the visible hand-off in the demo thread.")

# ============================================================
# SLIDE 13 — EXPIRING MEMBERS
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DAILY OPERATIONS", "Expiring members & renewal reminders")
text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.8),
     [[("The system finds members whose membership is about to lapse using a date-range query:", 16, DARK, True)]])
rounded(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.95), fill=NAVY)
text(s, Inches(0.9), Inches(2.62), Inches(11.5), Inches(0.7),
     [[("WHERE expiration_date BETWEEN CURRENT_DATE() AND DATE_ADD(CURRENT_DATE(), INTERVAL N DAY)",
        15, RGBColor(0x7E,0xE7,0x87), False, False, "Consolas")]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.6), Inches(3.85), Inches(12), Inches(3),
        [("Staff get a banner so they can call members before they expire.", 0),
         ("A member can be marked 'contacted' (renewal_contacted_at) so nobody is bothered twice.", 0),
         ("Backed by a composite index (expiration_date, renewal_contacted_at) for speed.", 0)],
        size=15, gap=14)
notes(s, "Great place to show the idx_expiration_date / composite index. Real, practical business value.")

# ============================================================
# SLIDE 14 — REFERENCE DATA
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "DAILY OPERATIONS", "Reference data: Plans · Promos · Expenses")
ref = [
    ("Plans", BLUE, "9 plans, ₱150 walk-in up to ₱12,000 / 12 months. Referenced by members & payments. ON DELETE RESTRICT means the DB refuses to delete a plan still in use."),
    ("Promos", ORANGE, "Bonus days or discounts (e.g. Happy Birthday Promo). is_active flag toggles availability without deleting."),
    ("Expenses", BLUE, "category, amount, date. Feeds the profit side of the reports in Part 4."),
]
y = Inches(1.75)
for i, (t, ac, d) in enumerate(ref):
    yy = y + i*Inches(1.65)
    rounded(s, Inches(0.6), yy, Inches(12.1), Inches(1.45), fill=LGRAY)
    box(s, Inches(0.6), yy, Inches(2.2), Inches(1.45), fill=ac)
    text(s, Inches(0.6), yy+Inches(0.45), Inches(2.2), Inches(0.6), [[(t, 19, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, Inches(3.1), yy+Inches(0.28), Inches(9.3), Inches(1.0), [[(d, 14.5, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Optional strong demo: try to delete a plan that has members -> the database rejects it. Integrity at the DB level, not just the UI.")

# ============================================================
# SLIDE 14b — CODE LOGIC (MEMBER 3)
# ============================================================
code_slide(
    "CODE LOGIC · PART 3", "The date query: get_expiring_members.php",
    [
        "SELECT",
        "    m.full_name,",
        "    m.expiration_date,",
        "    DATEDIFF(m.expiration_date, CURRENT_DATE())",
        "        AS days_left",
        "FROM members m",
        "LEFT JOIN plans pl",
        "    ON m.plan_id = pl.plan_id",
        "WHERE m.expiration_date >= CURRENT_DATE()",
        "  AND m.expiration_date <=",
        "      DATE_ADD(CURRENT_DATE(), INTERVAL 7 DAY)",
        "ORDER BY m.expiration_date ASC",
    ],
    [
        "DATEDIFF - how many days are left before the membership expires.",
        "WHERE ... from today up to +7 days -> only members about to lapse appear.",
        "LEFT JOIN plans - to pull each member's plan name and price.",
        "expiration_date is indexed -> stays fast even as members grow.",
    ],
    accent=BLUE, tag="Backend/get_expiring_members.php",
)

# ============================================================
# DIVIDER PART 4
# ============================================================
divider("4", "Gabriel Caporado", "Insight & Trust: Reporting + Security",
        "Turning rows into reports with aggregation SQL — and protecting the database from login to audit log.",
        accent=ORANGE)

# ============================================================
# SLIDE 15 — REPORTING
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "INSIGHT", "Reporting = aggregation queries")
text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.6),
     [[("The heavy lifting happens in SQL, not PHP. This is what a relational database is built for.", 16, DARK, True)]])
rounded(s, Inches(0.6), Inches(2.25), Inches(12.1), Inches(1.0), fill=NAVY)
text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.7),
     [[("SELECT SUM(amount), customer_type FROM payments JOIN plans ... GROUP BY period, customer_type",
        14, RGBColor(0x7E,0xE7,0x87), False, False, "Consolas")]], anchor=MSO_ANCHOR.MIDDLE)
rep = [
    ("Revenue report", "SUM() + GROUP BY by day / week / month, split by member vs walk-in."),
    ("Dashboard stats", "COUNT() of members, visits, active vs expiring."),
    ("Branch sales report", "Actual revenue vs monthly_targets, reconciled with bank_deposits (variance)."),
    ("JOINs everywhere", "payments → plans → members to label every row."),
]
y = Inches(3.6)
for i, (t, d) in enumerate(rep):
    col = i % 2; row = i // 2
    x = Inches(0.6)+col*Inches(6.2); yy = y+row*Inches(1.5)
    rounded(s, x, yy, Inches(5.95), Inches(1.35), fill=LGRAY)
    box(s, x, yy, Inches(0.1), Inches(1.35), fill=ORANGE)
    text(s, x+Inches(0.3), yy+Inches(0.16), Inches(5.5), Inches(1.1),
         [[(t, 15, NAVY, True)],[(d, 12.5, DARK, False)]], space_after=4)
notes(s, "Live: open Revenue report and show the payment from the member registered earlier now in the totals. Closes the demo loop.")

# ============================================================
# SLIDE 16 — SECURITY
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "TRUST", "Security: hardened, not just a password check")
sec = [
    ("bcrypt password hashing", "Passwords are hashed one-way — even we can't read them."),
    ("Token sessions", "Random 64-char token in the sessions table, 24-hour expiry, sent each request."),
    ("Rate limiting", "5 failed logins per IP in 15 min = locked. Tracked in login_attempts."),
    ("Timing-attack defense", "Fake bcrypt work when the username is wrong, so attackers can't enumerate usernames by response time."),
    ("Prepared statements", "Every query is parameterized — SQL-injection safe."),
    ("Session cleanup", "Expired sessions are deleted automatically."),
]
y = Inches(1.65)
for i, (t, d) in enumerate(sec):
    col = i % 2; row = i // 2
    x = Inches(0.6)+col*Inches(6.2); yy = y+row*Inches(1.7)
    rounded(s, x, yy, Inches(5.95), Inches(1.55), fill=LGRAY)
    box(s, x, yy, Inches(0.1), Inches(1.55), fill=BLUE if col==0 else ORANGE)
    text(s, x+Inches(0.32), yy+Inches(0.13), Inches(5.5), Inches(1.4),
         [[(t, 14.5, NAVY, True)],[(d, 11.5, DARK, False)]], space_after=4)
notes(s, "This is the most impressive area. Optionally demo a live lockout by typing a wrong password 5x.")

# ============================================================
# SLIDE 17 — AUDIT + INTEGRITY
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "TRUST", "Audit trail & data integrity")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.1), Inches(4.5),
        [("Every write is logged to activity_log:", 0),
         ("who (admin + username), when, their IP", 1),
         ("a JSON snapshot of exactly what changed", 1),
         ("Full accountability — see who changed what.", 0),
         ("Fail-safe: if logging breaks, it never breaks the real operation.", 0)],
        size=15, gap=12)
rounded(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.6), fill=LBLUE)
text(s, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.5), [[("Integrity enforced by the DB itself", 16, NAVY, True)]])
bullets(s, Inches(7.2), Inches(2.6), Inches(5.2), Inches(3.5),
        [("Foreign keys — no orphan payments, no deleting an in-use plan.", 0),
         ("NOT NULL & DEFAULT constraints on critical columns.", 0),
         ("Transactions — multi-table writes are all-or-nothing.", 0)],
        size=14, gap=14)
notes(s, "Show the Activity Log view with the member.create entry from the demo. Tie integrity back to the ERD foreign keys.")

# ============================================================
# SLIDE 17b — CODE LOGIC (MEMBER 4)
# ============================================================
code_slide(
    "CODE LOGIC · PART 4", "Aggregation SQL: get_revenue.php",
    [
        "SELECT",
        "    MONTH(payment_date)     AS month_num,",
        "    MONTHNAME(payment_date) AS month_name,",
        "    SUM(amount)             AS total,",
        "    COUNT(*)                AS count",
        "FROM payments",
        "WHERE YEAR(payment_date) = :year",
        "GROUP BY MONTH(payment_date)",
        "ORDER BY month_num ASC",
        "",
        "-- Example result:",
        "--  May  2026  |  total 70,650  |  count 19",
    ],
    [
        "SUM(amount) - total revenue; COUNT(*) - number of transactions.",
        "GROUP BY month -> a separate total per month, not lumped together.",
        "WHERE YEAR = :year -> only the chosen year (bound parameter).",
        "The DATABASE does the math, not PHP - this is the power of SQL.",
    ],
    accent=ORANGE, tag="Backend/get_revenue.php",
)

# ============================================================
# SLIDE 18 — DEMO FLOW
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "LIVE DEMO", "One continuous story")
demo = [
    ("Marco", "Trace a click through the 3 tiers", BLUE),
    ("Ariz", "Register a brand-new member, live", ORANGE),
    ("Brix", "Time that same member in", BLUE),
    ("Gabriel", "Show their payment in the revenue report + the audit log", ORANGE),
]
for i, (who, what, ac) in enumerate(demo):
    yy = Inches(1.8)+i*Inches(1.25)
    rounded(s, Inches(0.8), yy, Inches(11.7), Inches(1.05), fill=LGRAY)
    box(s, Inches(0.8), yy, Inches(0.12), Inches(1.05), fill=ac)
    text(s, Inches(1.2), yy+Inches(0.28), Inches(2.6), Inches(0.6), [[(who, 18, ac, True)]])
    text(s, Inches(4.0), yy+Inches(0.3), Inches(8.2), Inches(0.6), [[(what, 15.5, DARK, False)]])
    if i < 3:
        text(s, Inches(6.3), yy+Inches(1.0), Inches(1), Inches(0.3), [[("↓", 18, ORANGE, True)]], align=PP_ALIGN.CENTER)
notes(s, "Rehearse this end-to-end at least once. The same record visibly travels through all four presenters.")

# ============================================================
# SLIDE 19 — CONCLUSION
# ============================================================
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
box(s, 0, Inches(2.4), EMU_W, Inches(0.06), fill=ORANGE)
text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1),
     [[("Fitness Synergy, end to end", 40, WHITE, True)]])
text(s, Inches(0.9), Inches(2.7), Inches(11.6), Inches(3),
     [[("A normalized relational database (3NF) with foreign keys & indexes", 17, LBLUE, False)],
      [("A clean 3-tier app: React → PHP REST API → MySQL", 17, LBLUE, False)],
      [("Real transactions & constraints keeping data consistent", 17, LBLUE, False)],
      [("Security from login to audit log", 17, LBLUE, False)]], space_after=12)
text(s, Inches(0.9), Inches(6.0), Inches(11.6), Inches(1),
     [[("Thank you! — We're happy to take your questions.", 20, ORANGE, True)]])
notes(s, "Member 4 closes. Then open the floor. Everyone should know the ERD for Q&A.")

# ============================================================
# Stamp presenter name (first name) on each content slide.
# Title (1), dividers (6,11,16,21) and closing (27) are skipped —
# those already show the presenter prominently.
# ============================================================
pres_map = {
    2: "Marco", 3: "Marco", 4: "Marco", 5: "Marco",
    7: "Marco", 8: "Marco", 9: "Marco", 10: "Marco",
    11: "Marco", 12: "Marco", 13: "Marco", 14: "Marco",
    16: "Ariz", 17: "Ariz", 18: "Ariz", 19: "Ariz",
    21: "Brix", 22: "Brix", 23: "Brix", 24: "Brix",
    26: "Gabriel", 27: "Gabriel", 28: "Gabriel", 29: "Gabriel",
    30: "Whole Team",
}
acc_map = {"Marco": BLUE, "Ariz": ORANGE, "Brix": BLUE, "Gabriel": ORANGE, "Whole Team": NAVY}
for _i, _sl in enumerate(list(prs.slides), start=1):
    if _i not in pres_map:
        continue
    _name = pres_map[_i]; _ac = acc_map[_name]
    rounded(_sl, Inches(10.05), Inches(6.98), Inches(2.78), Inches(0.42), fill=LGRAY)
    box(_sl, Inches(10.05), Inches(6.98), Inches(0.1), Inches(0.42), fill=_ac)
    text(_sl, Inches(10.28), Inches(7.0), Inches(2.5), Inches(0.38),
         [[("Presenter:  ", 10, GRAY, False), (_name, 11, NAVY, True)]],
         anchor=MSO_ANCHOR.MIDDLE)

prs.save("Fitness_Synergy_Presentation.pptx")
print("Saved Fitness_Synergy_Presentation.pptx  —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
